"""
Generate PowerPoint presentation for UgaJapa internship demo.
Run: python docs/generate_presentation.py
Output: docs/UgaJapa-Presentation.pptx
"""

from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

DOCS_DIR = os.path.dirname(os.path.abspath(__file__))
SCREENSHOTS = os.path.join(DOCS_DIR, "screenshots")
OUTPUT = os.path.join(DOCS_DIR, "UgaJapa-Presentation.pptx")

# Modern palette — deep navy + indigo accent + soft surfaces
NAVY = RGBColor(15, 23, 42)
INDIGO = RGBColor(79, 70, 229)
INDIGO_LIGHT = RGBColor(224, 231, 255)
SKY = RGBColor(14, 165, 233)
EMERALD = RGBColor(16, 185, 129)
AMBER = RGBColor(245, 158, 11)
SLATE = RGBColor(51, 65, 85)
MUTED = RGBColor(100, 116, 139)
WHITE = RGBColor(255, 255, 255)
BG = RGBColor(248, 250, 252)
CARD = RGBColor(255, 255, 255)
BORDER = RGBColor(226, 232, 240)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
_slide_no = 0


def next_slide_no() -> int:
    global _slide_no
    _slide_no += 1
    return _slide_no


def blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_footer(slide)
    return slide


def set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_footer(slide, dark: bool = False) -> None:
    n = next_slide_no()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.12), SLIDE_W, Inches(0.38))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY if dark else BORDER
    bar.line.fill.background()

    label = slide.shapes.add_textbox(Inches(0.45), Inches(7.14), Inches(8), Inches(0.3))
    tf = label.text_frame
    tf.text = "UgaJapa Translation for Mattermost · v2.3.0"
    p = tf.paragraphs[0]
    p.font.size = Pt(9)
    p.font.color.rgb = WHITE if dark else MUTED

    num = slide.shapes.add_textbox(Inches(12.35), Inches(7.14), Inches(0.7), Inches(0.3))
    ntf = num.text_frame
    ntf.text = str(n)
    np = ntf.paragraphs[0]
    np.font.size = Pt(9)
    np.font.color.rgb = WHITE if dark else MUTED
    np.alignment = PP_ALIGN.RIGHT


def add_header(slide, title: str, subtitle: str | None = None) -> None:
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), Inches(1.05))
    accent.fill.solid()
    accent.fill.fore_color.rgb = INDIGO
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.42), Inches(0.22), Inches(12.5), Inches(0.55))
    ttf = title_box.text_frame
    ttf.text = title
    ttf.paragraphs[0].font.size = Pt(32)
    ttf.paragraphs[0].font.bold = True
    ttf.paragraphs[0].font.color.rgb = NAVY

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.42), Inches(0.72), Inches(12.5), Inches(0.35))
        stf = sub.text_frame
        stf.text = subtitle
        stf.paragraphs[0].font.size = Pt(14)
        stf.paragraphs[0].font.color.rgb = MUTED

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.42), Inches(1.08), Inches(12.5), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER
    line.line.fill.background()


def rounded_card(slide, left, top, width, height, fill=CARD, line=BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def text_in_shape(shape, text: str, size=14, bold=False, color=SLATE, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align


def add_numbered_items(slide, items: list[str], left=0.55, top=1.35, width=12.2, size=20):
    y = top
    for i, item in enumerate(items, start=1):
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, y, Inches(0.38), Inches(0.38))
        badge.fill.solid()
        badge.fill.fore_color.rgb = INDIGO
        badge.line.fill.background()
        text_in_shape(badge, str(i), size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        box = slide.shapes.add_textbox(left + Inches(0.52), y - Inches(0.02), width - Inches(0.55), Inches(0.55))
        tf = box.text_frame
        tf.word_wrap = True
        tf.text = item
        tf.paragraphs[0].font.size = Pt(size)
        tf.paragraphs[0].font.color.rgb = SLATE
        y += Inches(0.62)


def add_bullet_cards(slide, items: list[tuple[str, str]], cols: int = 2):
    """items: list of (title, description)"""
    card_w = Inches(5.95)
    card_h = Inches(1.05)
    gap_x = Inches(0.35)
    gap_y = Inches(0.22)
    start_x = Inches(0.55)
    start_y = Inches(1.35)

    for i, (title, desc) in enumerate(items):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        card = rounded_card(slide, x, y, card_w, card_h, fill=WHITE)
        tf = card.text_frame
        tf.clear()
        tf.margin_left = Pt(12)
        tf.margin_right = Pt(10)
        tf.margin_top = Pt(8)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = INDIGO
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = MUTED


def add_arch_box(slide, x, y, w, h, title: str, lines: list[str], fill, accent):
    card = rounded_card(slide, x, y, w, h, fill=fill)
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.08))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()

    tf = card.text_frame
    tf.clear()
    tf.margin_left = Pt(10)
    tf.margin_top = Pt(14)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = NAVY
    for line in lines:
        bp = tf.add_paragraph()
        bp.text = line
        bp.font.size = Pt(10)
        bp.font.color.rgb = SLATE
        bp.space_before = Pt(2)


def add_arrow_down(slide, x, y, h=Inches(0.35)):
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, y, Inches(0.28), h)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = INDIGO
    arrow.line.fill.background()


def add_arrow_right(slide, x, y, w=Inches(0.45)):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, Inches(0.22))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = INDIGO
    arrow.line.fill.background()


def add_architecture_slide(prs: Presentation):
    slide = blank_slide(prs)
    add_header(slide, "System Architecture", "Three layers — browser, Mattermost plugin, Translation API, Google Cloud")

    # Zone labels
    zones = [
        (Inches(0.45), Inches(1.22), Inches(3.55), "USER", SKY),
        (Inches(4.15), Inches(1.22), Inches(4.55), "MATTERMOST (Docker)", INDIGO),
        (Inches(8.85), Inches(1.22), Inches(4.0), "SERVICES", EMERALD),
    ]
    for x, y, w, label, color in zones:
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.28))
        pill.fill.solid()
        pill.fill.fore_color.rgb = color
        pill.line.fill.background()
        text_in_shape(pill, label, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Row 1 — main flow
    add_arch_box(
        slide, Inches(0.55), Inches(1.62), Inches(3.35), Inches(1.15),
        "Browser",
        ["Mattermost web UI", "Plugin React webapp", "WhatsApp-style chat"],
        INDIGO_LIGHT, INDIGO,
    )
    add_arrow_right(slide, Inches(3.98), Inches(2.05))
    add_arch_box(
        slide, Inches(4.55), Inches(1.62), Inches(4.15), Inches(1.15),
        "Mattermost 10.5",
        ["Team chat server", "PostgreSQL database", "Plugin server (Go)"],
        RGBColor(237, 233, 254), INDIGO,
    )
    add_arrow_right(slide, Inches(8.78), Inches(2.05))
    add_arch_box(
        slide, Inches(9.35), Inches(1.62), Inches(3.45), Inches(1.15),
        "Translation API",
        ["Node.js · port 5000", "Translate · STT · TTS", "Quality scoring"],
        RGBColor(209, 250, 229), EMERALD,
    )

    add_arrow_down(slide, Inches(6.52), Inches(2.85), Inches(0.4))

    # Row 2 — Google + fallback
    add_arch_box(
        slide, Inches(3.2), Inches(3.35), Inches(3.6), Inches(1.2),
        "Google Cloud (primary)",
        ["Cloud Translation API", "Speech-to-Text", "Text-to-Speech"],
        RGBColor(254, 243, 199), AMBER,
    )
    add_arch_box(
        slide, Inches(7.1), Inches(3.35), Inches(3.0), Inches(1.2),
        "Fallback",
        ["MyMemory (no Google key)", "Whisper (STT backup)"],
        RGBColor(241, 245, 249), MUTED,
    )

    # Legend strip — what we built
    legend = rounded_card(slide, Inches(0.55), Inches(4.85), Inches(12.25), Inches(1.55), fill=WHITE)
    tf = legend.text_frame
    tf.clear()
    tf.margin_left = Pt(16)
    tf.margin_top = Pt(10)
    p = tf.paragraphs[0]
    p.text = "What we built"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY

    parts = [
        "① mattermost-plugin-translation — UI, auto-translate, voice/video hooks",
        "② translation-api — Google integration, slang, quality scores",
        "③ docker-compose.yml — local Mattermost + PostgreSQL demo",
    ]
    for part in parts:
        bp = tf.add_paragraph()
        bp.text = part
        bp.font.size = Pt(12)
        bp.font.color.rgb = SLATE
        bp.space_before = Pt(4)

    # Flow note
    note = slide.shapes.add_textbox(Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.45))
    ntf = note.text_frame
    ntf.text = "Message flow: User sends text → Mattermost stores original → Plugin calls API per reader → Translation appears in each user's language"
    ntf.paragraphs[0].font.size = Pt(11)
    ntf.paragraphs[0].font.italic = True
    ntf.paragraphs[0].font.color.rgb = MUTED


def add_flow_pipeline(slide, steps: list[tuple[str, str]], top=Inches(1.4)):
    n = len(steps)
    card_w = Inches(11.8 / n - 0.12)
    gap = Inches(0.12)
    x = Inches(0.55)
    colors = [INDIGO, SKY, EMERALD, AMBER, INDIGO, SKY]

    for i, (num, label) in enumerate(steps):
        card = rounded_card(slide, x, top, card_w, Inches(1.35), fill=WHITE)
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.12), top + Inches(0.12), Inches(0.35), Inches(0.35))
        badge.fill.solid()
        badge.fill.fore_color.rgb = colors[i % len(colors)]
        badge.line.fill.background()
        text_in_shape(badge, num, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        tf = card.text_frame
        tf.clear()
        tf.margin_left = Pt(8)
        tf.margin_top = Pt(42)
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(10)
        p.font.color.rgb = SLATE
        p.alignment = PP_ALIGN.CENTER

        if i < n - 1:
            add_arrow_right(slide, x + card_w + Inches(0.02), top + Inches(0.55), Inches(0.1))
        x += card_w + gap


def add_image_slide(prs: Presentation, title: str, subtitle: str, image_name: str, callouts: list[str]):
    slide = blank_slide(prs)
    add_header(slide, title, subtitle)

    frame = rounded_card(slide, Inches(0.45), Inches(1.28), Inches(8.15), Inches(5.55), fill=WHITE)
    image_path = os.path.join(SCREENSHOTS, image_name)
    if os.path.isfile(image_path):
        slide.shapes.add_picture(image_path, Inches(0.58), Inches(1.42), width=Inches(7.9))
    else:
        text_in_shape(frame, f"Missing: {image_name}", size=14, color=MUTED, align=PP_ALIGN.CENTER)

    panel = rounded_card(slide, Inches(8.85), Inches(1.28), Inches(4.05), Inches(5.55), fill=INDIGO_LIGHT, line=INDIGO)
    tf = panel.text_frame
    tf.clear()
    tf.margin_left = Pt(14)
    tf.margin_top = Pt(12)
    hp = tf.paragraphs[0]
    hp.text = "Key points"
    hp.font.size = Pt(15)
    hp.font.bold = True
    hp.font.color.rgb = INDIGO

    for point in callouts:
        p = tf.add_paragraph()
        p.text = f"▸  {point}"
        p.font.size = Pt(12)
        p.font.color.rgb = SLATE
        p.space_before = Pt(10)


def add_title_slide(prs: Presentation):
    global _slide_no
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    _slide_no += 1

    # Decorative shapes
    c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(-0.8), Inches(3.5), Inches(3.5))
    c1.fill.solid()
    c1.fill.fore_color.rgb = INDIGO
    c1.fill.transparency = 0.35
    c1.line.fill.background()

    c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.2), Inches(5.2), Inches(2.8), Inches(2.8))
    c2.fill.solid()
    c2.fill.fore_color.rgb = EMERALD
    c2.fill.transparency = 0.55
    c2.line.fill.background()

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.55), Inches(1.1), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = EMERALD
    bar.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11.5), Inches(1.2))
    ttf = title.text_frame
    ttf.text = "UgaJapa Translation"
    ttf.paragraphs[0].font.size = Pt(46)
    ttf.paragraphs[0].font.bold = True
    ttf.paragraphs[0].font.color.rgb = WHITE

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(2.75), Inches(11.5), Inches(0.6))
    stf = sub.text_frame
    stf.text = "Real-time multilingual chat for Mattermost"
    stf.paragraphs[0].font.size = Pt(22)
    stf.paragraphs[0].font.color.rgb = RGBColor(203, 213, 225)

    chips = slide.shapes.add_textbox(Inches(0.8), Inches(3.55), Inches(11.5), Inches(1.5))
    ctf = chips.text_frame
    ctf.text = "Internship project  ·  v2.3.0  ·  June 2026"
    for line in [
        "",
        "Mahamat Abdelrassoul  ·  UgaJapa team",
        "Text · Voice · Video · Auto-translate · Read-aloud",
    ]:
        p = ctf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(148, 163, 184)

    add_footer(slide, dark=True)


def add_quality_measurement_slide(prs: Presentation):
    """TransChecker method — back-translation + 3 scores."""
    slide = blank_slide(prs)
    add_header(slide, "How We Measure Translation Quality", "TransChecker method — automatic, no human reviewer needed")

    # Back-translation flow
    steps = [
        ("1", "Original\ntext"),
        ("2", "Translate\nto target"),
        ("3", "Translate\nback"),
        ("4", "Compare\nmeaning"),
    ]
    add_flow_pipeline(slide, steps, top=Inches(1.32))

    note = slide.shapes.add_textbox(Inches(0.55), Inches(2.85), Inches(12.2), Inches(0.55))
    ntf = note.text_frame
    ntf.text = (
        "Example: French “Bonjour” → English “Hello” → back to French “Bonjour” → high match = good translation"
    )
    ntf.paragraphs[0].font.size = Pt(12)
    ntf.paragraphs[0].font.italic = True
    ntf.paragraphs[0].font.color.rgb = MUTED

    # Three score cards + combined
    scores = [
        ("Character match", "Levenshtein", "Are letters similar after back-translation?", INDIGO_LIGHT, INDIGO),
        ("Word overlap", "semantic_score", "Do the same important words appear?", RGBColor(224, 242, 254), SKY),
        ("AI meaning", "embedding_score", "MiniLM model — same meaning even if words differ?", RGBColor(209, 250, 229), EMERALD),
        ("Final score", "quality_score", "Combined % — API picks the best candidate", RGBColor(254, 243, 199), AMBER),
    ]
    x_positions = [Inches(0.55), Inches(3.35), Inches(6.15), Inches(8.95)]
    w = Inches(2.55)
    for (title, field, desc, fill, accent), x in zip(scores, x_positions):
        card = rounded_card(slide, x, Inches(3.55), w, Inches(2.35), fill=fill)
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(3.55), w, Inches(0.07))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = accent
        stripe.line.fill.background()
        tf = card.text_frame
        tf.clear()
        tf.margin_left = Pt(8)
        tf.margin_top = Pt(12)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p2 = tf.add_paragraph()
        p2.text = field
        p2.font.size = Pt(9)
        p2.font.color.rgb = accent
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(10)
        p3.font.color.rgb = SLATE

    foot = slide.shapes.add_textbox(Inches(0.55), Inches(6.15), Inches(12.2), Inches(0.85))
    ftf = foot.text_frame
    ftf.word_wrap = True
    ftf.text = "For tricky short messages, the API may try multiple paths and keep the highest quality_score."
    p = ftf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.color.rgb = MUTED
    p2 = ftf.add_paragraph()
    p2.text = "Voice & video use fast mode (one Google call) — no quality scoring, by design for speed."
    p2.font.size = Pt(11)
    p2.font.bold = True
    p2.font.color.rgb = INDIGO


def add_quality_visibility_slide(prs: Presentation):
    """Where supervisors / users can actually see the scores."""
    slide = blank_slide(prs)
    add_header(slide, "Where Can You See Quality?", "How to demonstrate scores live during the presentation")

    add_bullet_cards(
        slide,
        [
            ("/translate command", "Type: /translate Hello — shows translation + Quality % + AI semantic % + character match + back-translation"),
            ("Pre-send preview", "When you type in a foreign language, popup shows: AI quality 87% before you send"),
            ("API response", "POST /translate returns quality_score, score, semantic_score, embedding_score (for integrations)"),
            ("Normal chat messages", "No % shown on every bubble — keeps chat clean for daily use"),
            ("Voice & video", "No quality % — fast path only; show transcript translation instead"),
            ("Health check", "GET /health shows google-translate + semantic_embeddings_enabled"),
        ],
        cols=2,
    )

    demo = rounded_card(slide, Inches(0.55), Inches(5.35), Inches(12.25), Inches(1.35), fill=INDIGO_LIGHT, line=INDIGO)
    tf = demo.text_frame
    tf.clear()
    tf.margin_left = Pt(14)
    tf.margin_top = Pt(10)
    p = tf.paragraphs[0]
    p.text = "Best live demo for your audience tomorrow"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = INDIGO
    for line in [
        "1. In Mattermost chat, run:  /translate Bonjour tout le monde",
        "2. Show the reply with Quality 85% · AI semantic 90% · Character match 70%",
        "3. Explain: we translate → translate back → compare → combine into one percentage",
    ]:
        bp = tf.add_paragraph()
        bp.text = line
        bp.font.size = Pt(12)
        bp.font.color.rgb = SLATE
        bp.space_before = Pt(4)


def add_closing_slide(prs: Presentation):
    global _slide_no
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    _slide_no += 1

    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.8), Inches(0.5), Inches(3.8), Inches(3.8))
    c.fill.solid()
    c.fill.fore_color.rgb = INDIGO
    c.fill.transparency = 0.4
    c.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.8), Inches(2.8))
    tf = box.text_frame
    tf.text = "Thank you"
    tf.paragraphs[0].font.size = Pt(48)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "\nQuestions?"
    p2.font.size = Pt(26)
    p2.font.color.rgb = RGBColor(203, 213, 225)
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "\ngithub.com/MahamatAbdelrassoul/MATTERMOST_PLUGIN"
    p3.font.size = Pt(13)
    p3.font.color.rgb = EMERALD
    p3.alignment = PP_ALIGN.CENTER

    add_footer(slide, dark=True)


def build() -> None:
    global _slide_no
    _slide_no = 0

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    add_title_slide(prs)

    # Agenda — 2-column cards
    slide = blank_slide(prs)
    add_header(slide, "Agenda", "What we will cover in this presentation")
    add_bullet_cards(
        slide,
        [
            ("Problem", "Why multilingual teams need help in chat"),
            ("Live demo", "Same conversation in English and Japanese"),
            ("Features", "Text, voice, video, read-aloud, preview"),
            ("Quality scoring", "How we measure translation accuracy (TransChecker)"),
            ("Architecture", "How the plugin and API work together"),
            ("Tech stack", "Mattermost, Go, React, Node, Google Cloud"),
            ("Delivery", "What we built and how to run the demo"),
        ],
        cols=2,
    )

    # Problem — visual pain points
    slide = blank_slide(prs)
    add_header(slide, "The Problem", "International teams chat in one place — but not one language")
    add_bullet_cards(
        slide,
        [
            ("Language barrier", "Team members speak English, Japanese, French… in the same channel"),
            ("Broken workflow", "Copy-paste into Google Translate kills conversation flow"),
            ("Voice & video", "Audio messages are impossible to understand across languages"),
            ("No native tool", "Mattermost has no built-in per-user translation"),
            ("Our goal", "Chat naturally — each person reads in their own language"),
        ],
        cols=2,
    )

    # Solution
    slide = blank_slide(prs)
    add_header(slide, "Our Solution", "Plugin + API — no changes to Mattermost core")
    add_bullet_cards(
        slide,
        [
            ("Per-user language", "Each person picks their receive language in settings"),
            ("Auto-translate", "Readers see translations; senders keep original text"),
            ("Voice & video", "Record → transcribe → translate for every reader"),
            ("Modern UI", "WhatsApp-style bubbles + language badges"),
            ("Read-aloud", "Speaker icon uses Google Text-to-Speech"),
            ("Quality check", "TransChecker scores translation confidence (text)"),
        ],
        cols=2,
    )

    # Demo slides
    add_image_slide(
        prs,
        "Live Demo — English View",
        "User Abdel reads and writes in English",
        "01-abdel-english-chat.png",
        [
            "Green bubbles on the right = Abdel's messages",
            "White bubbles on the left = Yohana's messages",
            "Conversation feels like normal English chat",
            "Speaker icon reads any message aloud",
        ],
    )
    add_image_slide(
        prs,
        "Live Demo — Japanese View",
        "Same chat, same time — Yohana sees Japanese",
        "02-yohana-japanese-chat.png",
        [
            "This is the proof slide for your presentation",
            "Abdel typed in English — Yohana sees Japanese",
            "No manual re-typing or copy-paste",
            "Each user reads in their own language",
        ],
    )

    # Features grid
    slide = blank_slide(prs)
    add_header(slide, "Key Features", "Plugin v2.3.0 — everything in one chat experience")
    add_bullet_cards(
        slide,
        [
            ("Auto-translate", "Every new message translated per reader"),
            ("Pre-send preview", "See translation before sending foreign text"),
            ("Voice messages", "Google STT → translate → text under player"),
            ("Video messages", "Audio extracted, same pipeline as voice"),
            ("Languages panel", "Member badges: EN, JA, FR…"),
            ("Slash commands", "/translation-lang and /translate"),
            ("Chat slang", "bjr → bonjour before translating"),
            ("Usage tracking", "Monthly character count for billing"),
        ],
        cols=2,
    )

    # Architecture — visual diagram
    add_architecture_slide(prs)

    # Tech stack — colored rows
    slide = blank_slide(prs)
    add_header(slide, "Technology Stack", "Built on proven open-source and cloud tools")
    rows = [
        ("Chat platform", "Mattermost 10.5 · Docker · PostgreSQL 15", INDIGO_LIGHT, INDIGO),
        ("Plugin server", "Go · Mattermost Plugin API · WebSockets", RGBColor(224, 242, 254), SKY),
        ("Plugin UI", "React · TypeScript · Redux", RGBColor(237, 233, 254), INDIGO),
        ("Translation API", "Node.js · Express · TypeScript", RGBColor(209, 250, 229), EMERALD),
        ("Cloud AI", "Google Translate · Speech-to-Text · TTS", RGBColor(254, 243, 199), AMBER),
    ]
    y = Inches(1.35)
    for title, desc, fill, accent in rows:
        card = rounded_card(slide, Inches(0.55), y, Inches(12.25), Inches(0.95), fill=fill)
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), y, Inches(0.1), Inches(0.95))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = accent
        stripe.line.fill.background()
        tf = card.text_frame
        tf.clear()
        tf.margin_left = Pt(18)
        tf.margin_top = Pt(8)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = SLATE
        y += Inches(1.05)

    # Text translation flow — pipeline
    slide = blank_slide(prs)
    add_header(slide, "How Text Translation Works", "Simple 6-step pipeline — easy to explain to anyone")
    add_flow_pipeline(
        slide,
        [
            ("1", "User\nsends message"),
            ("2", "Mattermost\nstores original"),
            ("3", "Plugin checks\neach reader's language"),
            ("4", "API calls\nGoogle Translate"),
            ("5", "Quality score\n(text only)"),
            ("6", "Reader sees\ntranslation"),
        ],
    )
    note = slide.shapes.add_textbox(Inches(0.55), Inches(3.0), Inches(12.2), Inches(3.5))
    ntf = note.text_frame
    ntf.word_wrap = True
    ntf.text = "Important rule for the audience:"
    ntf.paragraphs[0].font.size = Pt(16)
    ntf.paragraphs[0].font.bold = True
    ntf.paragraphs[0].font.color.rgb = INDIGO
    for line in [
        "The sender always sees their own original words.",
        "Only readers with a different receive language see a translation.",
        "Example: Abdel writes English → Yohana reads Japanese automatically.",
    ]:
        p = ntf.add_paragraph()
        p.text = line
        p.font.size = Pt(15)
        p.font.color.rgb = SLATE
        p.space_before = Pt(8)

    # Quality measurement — TransChecker (likely Q&A topic)
    add_quality_measurement_slide(prs)
    add_quality_visibility_slide(prs)

    # Voice + screenshot
    add_image_slide(
        prs,
        "Voice & Video Messages",
        "Microphone and camera buttons in the composer",
        "04-voice-video-buttons.png",
        [
            "User records voice or short video",
            "File uploaded to Mattermost",
            "Server runs Google Speech-to-Text",
            "Fast translate path for speed",
            "Reader sees transcript in their language",
        ],
    )

    # Language panel
    add_image_slide(
        prs,
        "Language Settings",
        "Each user controls their own experience",
        "03-translation-languages-panel.png",
        [
            "Pick receive language (EN, JA, FR…)",
            "Choose read-aloud voice style",
            "See all channel members' languages",
            "Open from translate icon in header",
        ],
    )

    # Challenges
    slide = blank_slide(prs)
    add_header(slide, "Challenges We Solved", "Real bugs we fixed during development")
    add_numbered_items(
        slide,
        [
            "Voice showed wrong language → fixed with server-side STT + detection (v2.3.0)",
            "Slow voice processing → fast media path (one Google call)",
            "API fetch errors → retries on plugin, API, and Google",
            "Pre-send preview missing → fixed mixed-language typing logic",
            "Chat UX → WhatsApp-style layout + language badges",
        ],
        top=Inches(1.35),
        size=18,
    )

    # Deliverables
    slide = blank_slide(prs)
    add_header(slide, "What We Delivered", "Complete internship deliverable")
    add_bullet_cards(
        slide,
        [
            ("Plugin v2.3.0", "com.transchecker.translation — upload to Mattermost"),
            ("Translation API", "Standalone Node service with Google Cloud"),
            ("Documentation", "README + full Word guide for non-programmers"),
            ("GitHub repo", "Public repository with setup instructions"),
            ("Working demo", "Two users, two languages, one channel"),
            ("Quality scoring", "TransChecker prototype for confidence"),
        ],
        cols=2,
    )

    # How to run
    slide = blank_slide(prs)
    add_header(slide, "Running the Demo", "5 steps to reproduce locally")
    add_numbered_items(
        slide,
        [
            "Start Translation API — npm install, .env, npm run dev (port 5000)",
            "Start Mattermost — docker compose up -d → localhost:8065",
            "Build plugin — build-bundle.ps1 → upload .tar.gz in System Console",
            "Configure plugin — API URL: http://host.docker.internal:5000",
            "Set receive languages → open two browsers → start chatting",
        ],
        top=Inches(1.35),
        size=19,
    )

    # Future
    slide = blank_slide(prs)
    add_header(slide, "Future Work", "Next steps after the internship")
    add_bullet_cards(
        slide,
        [
            ("Production deploy", "Server, HTTPS, Dockerize the Translation API"),
            ("GitHub Release", "Ready-to-download plugin .tar.gz bundle"),
            ("Security", "Strong API keys and Google billing controls"),
            ("UX polish", "Speaking-language picker before recording"),
        ],
        cols=2,
    )

    add_closing_slide(prs)

    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT} ({_slide_no} slides)")


if __name__ == "__main__":
    build()
